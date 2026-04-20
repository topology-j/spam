import sys
if sys.stdout and hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
import os
from dotenv import load_dotenv
load_dotenv(os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env"), override=True)

from fastapi import FastAPI, HTTPException, Depends, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel
from passlib.context import CryptContext
from jose import jwt, JWTError
import sqlite3
from datetime import datetime, timedelta
from typing import Optional, List
import os
import io
import csv
import uuid
import re
from contextlib import contextmanager
from contextvars import ContextVar
from datetime import datetime, timedelta, timezone
from langsmith import traceable
from langsmith import Client as LangSmithClient

# ?????????????????????????????
# APP
# ?????????????????????????????
app = FastAPI()


@app.on_event("startup")
def on_startup():
    try:
        print(
            "[LangSmith][startup] "
            f"TRACING={os.getenv('LANGSMITH_TRACING') or os.getenv('LANGCHAIN_TRACING_V2')} "
            f"PROJECT={os.getenv('LANGSMITH_PROJECT') or os.getenv('LANGCHAIN_PROJECT')} "
            f"ENDPOINT={os.getenv('LANGSMITH_ENDPOINT') or os.getenv('LANGCHAIN_ENDPOINT') or 'https://api.smith.langchain.com'} "
            f"KEY={(os.getenv('LANGSMITH_API_KEY') or os.getenv('LANGCHAIN_API_KEY') or '')[:15]}..."
        )
        from rag_engine import _load_vectorstore
        _load_vectorstore()
    except Exception:
        pass


app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ?????????????????????????????
# CONFIG
# ?????????????????????????????
SECRET_KEY = os.getenv("SECRET_KEY", "dev-secret-key")
ALGORITHM = "HS256"
TOKEN_EXPIRE_HOURS = 24
ALLOW_RAG_DIRECT_LEARN = os.getenv("ALLOW_RAG_DIRECT_LEARN", "false").strip().lower() == "true"
ALLOW_REPORT_VECTORSTORE_LEARN = os.getenv("ALLOW_REPORT_VECTORSTORE_LEARN", "false").strip().lower() == "true"

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
security = HTTPBearer()
_langsmith_client = None


def get_langsmith_client():
    global _langsmith_client
    if _langsmith_client is None:
        api_key = os.getenv("LANGSMITH_API_KEY") or os.getenv("LANGCHAIN_API_KEY")
        api_url = os.getenv("LANGSMITH_ENDPOINT") or os.getenv("LANGCHAIN_ENDPOINT") or "https://api.smith.langchain.com"
        if api_key:
            _langsmith_client = LangSmithClient(api_key=api_key, api_url=api_url)
    return _langsmith_client


def log_input_to_langsmith(run_name: str, text: str, username: str):
    client = get_langsmith_client()
    if client is None:
        print(f"[LangSmith][manual] client unavailable, skipping raw input logging for {run_name}")
        return
    run_id = uuid.uuid4()
    project_name = os.getenv("LANGSMITH_PROJECT") or os.getenv("LANGCHAIN_PROJECT") or "default"
    now = datetime.now(timezone.utc)
    client.create_run(
        name=run_name,
        run_type="chain",
        project_name=project_name,
        inputs={"text": text, "username": username},
        id=run_id,
        start_time=now,
    )
    client.update_run(
        run_id,
        end_time=now,
        outputs={"text": text, "username": username},
    )
    print(f"[LangSmith][manual] logged {run_name} text={text!r} username={username}")


def log_chat_input_to_langsmith(text: str, username: str):
    log_input_to_langsmith("chat_input_raw", text, username)


def log_report_input_to_langsmith(text: str, username: str):
    log_input_to_langsmith("report_input_raw", text, username)


def _ensure_learning_enabled(flag: bool, detail: str):
    if not flag:
        raise HTTPException(status_code=403, detail=detail)

# ?????????????????????????????
# DB
# ?????????????????????????????
DB_PATH = os.path.join(os.path.dirname(__file__), "spam_chat.db")


def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def load_dataset_rows(split: str, max_samples: Optional[int] = None):
    base_dir = os.path.dirname(__file__)
    split_files = {
        "train": os.path.join(base_dir, "..", "train_split.csv"),
        "val": os.path.join(base_dir, "..", "val_split.csv"),
        "test": os.path.join(base_dir, "..", "test.csv"),
    }

    dataset_path = split_files.get(split)
    if not dataset_path:
        raise HTTPException(status_code=400, detail="split은 'train', 'val', 'test' 중 하나여야 합니다.")
    if not os.path.exists(dataset_path):
        raise HTTPException(status_code=404, detail=f"{split} 데이터 파일이 없습니다: {os.path.basename(dataset_path)}")

    with open(dataset_path, "r", encoding="utf-8-sig", newline="") as f:
        reader = csv.DictReader(f)
        rows = [
            {"text": row["text"], "label": row["label"].strip().lower()}
            for row in reader
            if row.get("text") and row.get("label")
        ]

    if max_samples is not None:
        rows = rows[:max_samples]
    return rows


def _resolve_dataset_path(split: str) -> str:
    base_dir = os.path.dirname(__file__)
    split_files = {
        "train": os.path.join(base_dir, "..", "train_split.csv"),
        "val": os.path.join(base_dir, "..", "val_split.csv"),
        "test": os.path.join(base_dir, "..", "test.csv"),
    }
    dataset_path = split_files.get(split)
    if not dataset_path:
        raise HTTPException(status_code=400, detail="split은 'train', 'val', 'test' 중 하나여야 합니다.")
    if not os.path.exists(dataset_path):
        raise HTTPException(status_code=404, detail=f"{split} 데이터 파일이 없습니다: {os.path.basename(dataset_path)}")
    return os.path.abspath(dataset_path)


def _load_eval_dataset(split: str, max_samples: Optional[int] = None):
    import pandas as pd

    eval_csv_path = _resolve_dataset_path(split)
    df = pd.read_csv(eval_csv_path, encoding="utf-8-sig").dropna(subset=["text", "label"])
    df["label"] = df["label"].astype(str).str.strip().str.lower()

    effective_df = df.head(max_samples).copy() if max_samples is not None else df.copy()
    rows = [
        {"text": row["text"], "label": row["label"]}
        for row in effective_df.to_dict(orient="records")
        if row.get("text") and row.get("label")
    ]
    return {
        "rows": rows,
        "eval_csv_path": eval_csv_path,
        "row_count": int(len(df)),
        "effective_row_count": len(rows),
        "first_3_titles": df["title"].head(3).tolist() if "title" in df.columns else "NO TITLE COL",
    }


def _log_eval_dataset(mode: str, split: str, eval_csv_path: str, row_count: int, first_3_titles, effective_row_count: int, model_path: str):
    import pandas as pd
    try:
        df_tmp = pd.read_csv(eval_csv_path, encoding="utf-8-sig").dropna(subset=["text", "label"])
        label_count = df_tmp["label"].str.strip().str.lower().value_counts().to_dict()
    except Exception as e:
        label_count = f"ERROR: {e}"
    print("==== DEBUG START ====")
    print("MODE:", mode)
    print("SPLIT:", split)
    print("FILE PATH:", eval_csv_path)
    print("ROW COUNT:", row_count)
    print("LABEL COUNT:", label_count)
    print("=====================")


def _hangul_ratio(text: str) -> float:
    if not text:
        return 0.0
    hangul = sum(1 for ch in text if "\uac00" <= ch <= "\ud7a3")
    return hangul / max(len(text), 1)


def _score_decoded_text(text: str) -> tuple[float, int]:
    if not text:
        return (-1.0, 0)
    cleaned = text.strip()
    if not cleaned:
        return (-0.5, 0)
    replacement_penalty = cleaned.count("\ufffd") * 3
    mojibake_penalty = cleaned.count("Ã") + cleaned.count("Â") + cleaned.count("ì") * 0.1
    hangul_bonus = _hangul_ratio(cleaned) * 100
    printable_bonus = sum(1 for ch in cleaned if ch.isprintable()) / max(len(cleaned), 1) * 10
    score = hangul_bonus + printable_bonus - replacement_penalty - mojibake_penalty
    return (score, len(cleaned))


def _decode_bytes_best_effort(content: bytes) -> str:
    candidates = []
    for enc in ("utf-8", "utf-8-sig", "cp949", "euc-kr", "iso-8859-1", "latin-1"):
        try:
            decoded = content.decode(enc)
            candidates.append((_score_decoded_text(decoded), decoded))
        except Exception:
            continue
    if not candidates:
        return content.decode("utf-8", errors="replace")
    candidates.sort(key=lambda item: item[0], reverse=True)
    return candidates[0][1]


def _repair_mojibake_text(text: str) -> str:
    candidates = [text]
    for src_enc, dst_enc in (
        ("latin-1", "utf-8"),
        ("cp1252", "utf-8"),
        ("latin-1", "cp949"),
        ("cp1252", "cp949"),
    ):
        try:
            repaired = text.encode(src_enc, errors="ignore").decode(dst_enc, errors="ignore")
            if repaired:
                candidates.append(repaired)
        except Exception:
            continue

    candidates = list(dict.fromkeys(candidates))
    scored = [(_score_decoded_text(candidate), candidate) for candidate in candidates]
    scored.sort(key=lambda item: item[0], reverse=True)
    return scored[0][1]


def _strip_html_tags(text: str) -> str:
    text = re.sub(r"<(script|style)\b[^>]*>.*?</\1>", " ", text, flags=re.IGNORECASE | re.DOTALL)
    text = re.sub(r"<br\s*/?>", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"</p\s*>", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    return text


def _normalize_extracted_text(text: str) -> str:
    text = _repair_mojibake_text(text)
    text = text.replace("\x00", " ")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_text_from_eml(content: bytes) -> str:
    from email import policy
    from email.header import decode_header, make_header
    from email.parser import BytesParser

    msg = BytesParser(policy=policy.default).parsebytes(content)
    parts = []

    subject = msg.get("subject")
    if subject:
        try:
            decoded_subject = str(make_header(decode_header(subject)))
        except Exception:
            decoded_subject = subject
        parts.append(f"제목: {decoded_subject}")

    def decode_part(part) -> str:
        payload = part.get_payload(decode=True)
        if payload is None:
            payload = (part.get_payload() or "").encode("utf-8", errors="ignore")
        charset = part.get_content_charset()
        if charset:
            try:
                return payload.decode(charset)
            except Exception:
                pass
        return _decode_bytes_best_effort(payload)

    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_maintype() == "multipart":
                continue
            content_type = part.get_content_type()
            disposition = (part.get("Content-Disposition") or "").lower()
            if "attachment" in disposition:
                continue
            body = decode_part(part)
            if content_type == "text/html":
                body = _strip_html_tags(body)
            parts.append(body)
    else:
        body = decode_part(msg)
        if msg.get_content_type() == "text/html":
            body = _strip_html_tags(body)
        parts.append(body)

    return _normalize_extracted_text("\n".join(p for p in parts if p))


def _extract_text_from_hwp_preview(content: bytes) -> str:
    import olefile

    with olefile.OleFileIO(io.BytesIO(content)) as ole:
        if not ole.exists("PrvText"):
            return ""
        raw = ole.openstream("PrvText").read()

    for enc in ("utf-16", "utf-16-le", "cp949", "euc-kr"):
        try:
            text = raw.decode(enc)
            if text.strip():
                return _normalize_extracted_text(text)
        except Exception:
            continue
    return ""


def _extract_text_from_hwp_body(content: bytes) -> str:
    import olefile
    import zlib

    candidates = []
    with olefile.OleFileIO(io.BytesIO(content)) as ole:
        section_paths = [
            entry for entry in ole.listdir()
            if len(entry) == 2 and entry[0] == "BodyText" and entry[1].startswith("Section")
        ]
        section_paths.sort(key=lambda entry: entry[1])

        for entry in section_paths:
            raw = ole.openstream(entry).read()
            byte_candidates = [raw]
            for wbits in (-15, zlib.MAX_WBITS):
                try:
                    byte_candidates.append(zlib.decompress(raw, wbits))
                except Exception:
                    continue

            for payload in byte_candidates:
                for enc in ("utf-16-le", "utf-16", "cp949", "euc-kr"):
                    try:
                        decoded = payload.decode(enc, errors="ignore")
                        normalized = _normalize_extracted_text(decoded)
                        if normalized:
                            candidates.append(normalized)
                    except Exception:
                        continue

    if not candidates:
        return ""

    scored = [(_score_decoded_text(text), text) for text in candidates]
    scored.sort(key=lambda item: item[0], reverse=True)
    return scored[0][1]


def init_db():
    conn = get_db()
    cur = conn.cursor()
    cur.executescript("""
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        password_hash TEXT NOT NULL,
        nickname TEXT,
        name TEXT,
        phone TEXT,
        email TEXT,
        address TEXT,
        detail_address TEXT,
        postal_code TEXT,
        role TEXT NOT NULL DEFAULT 'user',
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    );

    CREATE TABLE IF NOT EXISTS spam_reports (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        email_content TEXT NOT NULL,
        status TEXT NOT NULL DEFAULT 'pending',
        counselor_id INTEGER,
        counselor_note TEXT DEFAULT '',
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (user_id) REFERENCES users(id),
        FOREIGN KEY (counselor_id) REFERENCES users(id)
    );

    CREATE TABLE IF NOT EXISTS spam_keywords (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        keyword TEXT UNIQUE NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    );

    CREATE TABLE IF NOT EXISTS chat_logs (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        username TEXT NOT NULL,
        role TEXT NOT NULL DEFAULT 'user',
        message TEXT NOT NULL,
        is_spam INTEGER,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (user_id) REFERENCES users(id)
    );

    CREATE TABLE IF NOT EXISTS counselor_reviews (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        report_id INTEGER NOT NULL UNIQUE,
        reviewer_id INTEGER NOT NULL,
        counselor_id INTEGER,
        stars INTEGER NOT NULL,
        comment TEXT DEFAULT '',
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (report_id) REFERENCES spam_reports(id),
        FOREIGN KEY (reviewer_id) REFERENCES users(id),
        FOREIGN KEY (counselor_id) REFERENCES users(id)
    );

    CREATE TABLE IF NOT EXISTS user_spam_rules (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        rule_type TEXT NOT NULL,
        value TEXT NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (user_id) REFERENCES users(id)
    );

    CREATE TABLE IF NOT EXISTS system_settings (
        key TEXT PRIMARY KEY,
        value TEXT NOT NULL
    );

    CREATE TABLE IF NOT EXISTS training_corrections (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        text TEXT NOT NULL,
        label TEXT NOT NULL,
        source TEXT DEFAULT 'improve',
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    );

    CREATE TABLE IF NOT EXISTS improvement_history (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        iteration INTEGER NOT NULL,
        mode TEXT NOT NULL DEFAULT 'fast',
        accuracy REAL, precision_score REAL, recall REAL, f1 REAL,
        total INTEGER, tp INTEGER, tn INTEGER, fp INTEGER, fn INTEGER,
        added_count INTEGER DEFAULT 0,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    );
    """)
    conn.commit()
    # detail_address 而щ읆 留덉씠洹몃젅?댁뀡
    try:
        conn.execute("ALTER TABLE users ADD COLUMN detail_address TEXT")
        conn.commit()
    except Exception:
        pass
    # 由щ럭 移댄뀒怨좊━ 蹂꾩젏 留덉씠洹몃젅?댁뀡
    for col in ["accuracy_stars", "processing_stars", "clarity_stars", "speed_stars"]:
        try:
            conn.execute(f"ALTER TABLE counselor_reviews ADD COLUMN {col} INTEGER")
            conn.commit()
        except Exception:
            pass
    conn.close()


def seed():
    conn = get_db()
    cur = conn.cursor()
    users = [
        ("admin", "admin123", "admin", "관리자", "관리자", "010-0000-0000", "admin@spam.com", "서울시", "00000"),
        ("counselor1", "counselor123", "counselor", "상담사1", "상담사", "010-1111-1111", "counselor1@spam.com", "서울시", "11111"),
        ("user1", "user123", "user", "사용자1", "사용자", "010-2222-2222", "user1@spam.com", "서울시", "22222"),
        ("green314", "1234", "developer", "그린", "개발자", "010-3333-3333", "green@spam.com", "서울시", "33333"),
    ]
    for username, pw, role, nickname, name, phone, email, address, postal_code in users:
        cur.execute("SELECT id FROM users WHERE username=?", (username,))
        if not cur.fetchone():
            cur.execute(
                "INSERT INTO users (username, password_hash, nickname, name, phone, email, address, postal_code, role) VALUES (?,?,?,?,?,?,?,?,?)",
                (username, pwd_context.hash(pw), nickname, name, phone, email, address, postal_code, role)
            )

    # 기본 스팸 키워드
    keywords = ["무료", "당첨", "클릭", "긴급", "송금", "비밀번호", "계좌번호", "투자", "광고"]
    for kw in keywords:
        cur.execute("INSERT OR IGNORE INTO spam_keywords (keyword) VALUES (?)", (kw,))

    # 湲곕낯 ?쒖뒪???ㅼ젙
    cur.execute("INSERT OR REPLACE INTO system_settings (key, value) VALUES ('report_enabled', 'true')")

    conn.commit()
    conn.close()


init_db()
seed()

# ?????????????????????????????
# AUTH
# ?????????????????????????????
def create_token(user_id, username, role):
    exp = datetime.utcnow() + timedelta(hours=TOKEN_EXPIRE_HOURS)
    return jwt.encode(
        {"sub": str(user_id), "username": username, "role": role, "exp": exp},
        SECRET_KEY, algorithm=ALGORITHM
    )


def get_current_user(cred: HTTPAuthorizationCredentials = Depends(security)):
    try:
        payload = jwt.decode(cred.credentials, SECRET_KEY, algorithms=[ALGORITHM])
        return {"id": int(payload["sub"]), "username": payload["username"], "role": payload["role"]}
    except JWTError:
        raise HTTPException(status_code=401, detail="?몄쬆 ?ㅻ쪟")


def require_role(*roles):
    def checker(user=Depends(get_current_user)):
        original = user["role"]
        effective = "admin" if original == "developer" else original
        if effective not in roles and original not in roles:
            raise HTTPException(status_code=403, detail="沅뚰븳 ?놁쓬")
        return user
    return checker


# ?????????????????????????????
# RAG ?붿쭊 濡쒕뱶 (GPT / Qwen)
# ?????????????????????????????
_rag_add_to_vectorstore = None
try:
    from rag_engine import classify_with_gpt, classify_with_qwen, add_to_vectorstore as _rag_add_to_vectorstore
    _rag_available = True
    print("??RAG ?붿쭊 濡쒕뱶 ?꾨즺 (GPT + Qwen)")
except Exception as e:
    print(f"?좑툘  RAG ?붿쭊 濡쒕뱶 ?ㅽ뙣: {e}")
    _rag_available = False


_eval_write_guard = ContextVar("eval_write_guard", default=False)


@contextmanager
def _evaluation_write_blocked():
    token = _eval_write_guard.set(True)
    try:
        yield
    finally:
        _eval_write_guard.reset(token)


def add_to_vectorstore(text: str, label: str) -> bool:
    if _eval_write_guard.get():
        raise RuntimeError("add_to_vectorstore() is blocked during evaluation paths")
    if _rag_add_to_vectorstore is None:
        raise RuntimeError("RAG vectorstore backend is not available")
    return _rag_add_to_vectorstore(text, label)


# ?????????????????????????????
# ?ㅽ뙵 媛먯? (GPT ?곗꽑 ??Qwen ???ㅼ썙???대갚)
# ?????????????????????????????
def detect_spam(text: str) -> bool:
    if _rag_available:
        # GPT濡?1李??먯젙
        try:
            result = classify_with_gpt(text)
            if "error" not in result:
                return result["is_spam"]
        except Exception as e:
            print(f"?좑툘  GPT ?먯젙 ?ㅽ뙣: {e}")

        # Qwen?쇰줈 2李??먯젙
        try:
            result = classify_with_qwen(text)
            if "error" not in result:
                return result["is_spam"]
        except Exception as e:
            print(f"?좑툘  Qwen ?먯젙 ?ㅽ뙣: {e}")

    # ?ㅼ썙???대갚
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT keyword FROM spam_keywords")
    keywords = [row["keyword"] for row in cur.fetchall()]
    conn.close()
    text_lower = text.lower()
    return any(kw.lower() in text_lower for kw in keywords)


# ?????????????????????????????
# LOGIN / REGISTER
# ?????????????????????????????
class LoginReq(BaseModel):
    username: str
    password: str


class RegisterReq(BaseModel):
    username: str
    password: str
    nickname: str
    name: str
    phone: str
    email: str
    address: str
    detail_address: Optional[str] = ""
    postal_code: str


@app.post("/auth/login")
def login(req: LoginReq):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT * FROM users WHERE username=?", (req.username,))
    user = cur.fetchone()
    conn.close()
    if not user or not pwd_context.verify(req.password, user["password_hash"]):
        raise HTTPException(status_code=401, detail="아이디 또는 비밀번호가 맞지 않습니다")
    return {"token": create_token(user["id"], user["username"], user["role"]), "username": user["username"], "role": user["role"]}


@app.post("/auth/register")
def register(req: RegisterReq):
    if len(req.username) < 2 or len(req.username) > 30:
        raise HTTPException(status_code=400, detail="아이디는 2~30자여야 합니다")
    if len(req.password) < 4:
        raise HTTPException(status_code=400, detail="비밀번호는 4자 이상이어야 합니다")
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT id FROM users WHERE username=?", (req.username,))
    if cur.fetchone():
        conn.close()
        raise HTTPException(status_code=400, detail="?대? ?ъ슜 以묒씤 ?꾩씠?붿엯?덈떎")
    cur.execute(
        "INSERT INTO users (username, password_hash, nickname, name, phone, email, address, detail_address, postal_code, role) VALUES (?,?,?,?,?,?,?,?,?,?)",
        (req.username, pwd_context.hash(req.password), req.nickname, req.name, req.phone, req.email, req.address, req.detail_address, req.postal_code, "user")
    )
    conn.commit()
    conn.close()
    return {"message": "?뚯썝媛???꾨즺"}


# ?????????????????????????????
# USERS
# ?????????????????????????????
@app.get("/users/me")
def get_me(user=Depends(get_current_user)):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT id, username, nickname, name, phone, email, address, detail_address, postal_code, role, created_at FROM users WHERE id=?", (user["id"],))
    row = cur.fetchone()
    conn.close()
    if not row:
        raise HTTPException(status_code=404, detail="?좎? ?놁쓬")
    return dict(row)


class ProfileUpdateReq(BaseModel):
    nickname: Optional[str] = None
    name: Optional[str] = None
    phone: Optional[str] = None
    email: Optional[str] = None
    address: Optional[str] = None
    detail_address: Optional[str] = None
    postal_code: Optional[str] = None
    current_password: Optional[str] = None
    new_password: Optional[str] = None


@app.patch("/users/me")
def update_me(req: ProfileUpdateReq, user=Depends(get_current_user)):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT * FROM users WHERE id=?", (user["id"],))
    row = cur.fetchone()
    if not row:
        conn.close()
        raise HTTPException(status_code=404, detail="?좎? ?놁쓬")

    new_hash = row["password_hash"]
    if req.new_password:
        if not req.current_password or not pwd_context.verify(req.current_password, row["password_hash"]):
            conn.close()
            raise HTTPException(status_code=400, detail="?꾩옱 鍮꾨?踰덊샇媛 ?щ컮瑜댁? ?딆뒿?덈떎")
        new_hash = pwd_context.hash(req.new_password)

    cur.execute("""
        UPDATE users SET
            nickname = COALESCE(?, nickname),
            name = COALESCE(?, name),
            phone = COALESCE(?, phone),
            email = COALESCE(?, email),
            address = COALESCE(?, address),
            detail_address = COALESCE(?, detail_address),
            postal_code = COALESCE(?, postal_code),
            password_hash = ?
        WHERE id = ?
    """, (req.nickname, req.name, req.phone, req.email, req.address, req.detail_address, req.postal_code, new_hash, user["id"]))
    conn.commit()
    conn.close()
    return {"message": "?섏젙 ?꾨즺"}


@app.get("/admin/users")
def get_users(user=Depends(require_role("admin"))):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT id, username, nickname, name, phone, email, address, postal_code, role, created_at FROM users ORDER BY created_at DESC")
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


class RoleUpdateReq(BaseModel):
    role: str


@app.patch("/admin/users/{user_id}/role")
def update_user_role(user_id: int, req: RoleUpdateReq, user=Depends(require_role("admin"))):
    caller_role = user["role"]
    if caller_role == "developer":
        allowed = {"user", "counselor", "admin", "developer"}
    else:
        allowed = {"user", "counselor", "admin"}
    if req.role not in allowed:
        raise HTTPException(status_code=403, detail="?대떦 ??븷濡?蹂寃쏀븷 沅뚰븳???놁뒿?덈떎")
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT id FROM users WHERE id=?", (user_id,))
    if not cur.fetchone():
        conn.close()
        raise HTTPException(status_code=404, detail="?ъ슜?먮? 李얠쓣 ???놁뒿?덈떎")
    cur.execute("UPDATE users SET role=? WHERE id=?", (req.role, user_id))
    conn.commit()
    conn.close()
    return {"message": "??븷??蹂寃쎈릺?덉뒿?덈떎"}


# ?????????????????????????????
# SPAM REPORTS
# ?????????????????????????????
class ReportReq(BaseModel):
    email_content: str


class ReportUpdateReq(BaseModel):
    status: str
    counselor_note: Optional[str] = ""
    keywords: Optional[List[str]] = []


@app.get("/spam-reports")
def get_reports(user=Depends(require_role("admin", "counselor"))):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("""
        SELECT r.*, u.username as requester
        FROM spam_reports r
        JOIN users u ON r.user_id = u.id
        ORDER BY r.created_at DESC
    """)
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


@app.get("/spam-reports/my")
def get_my_reports(user=Depends(get_current_user)):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("""
        SELECT r.*, rv.id as review_id, rv.stars as review_stars, rv.comment as review_comment,
               u.nickname as counselor_nickname, u.username as counselor_username
        FROM spam_reports r
        LEFT JOIN counselor_reviews rv ON rv.report_id = r.id
        LEFT JOIN users u ON r.counselor_id = u.id
        WHERE r.user_id = ?
        ORDER BY r.created_at DESC
    """, (user["id"],))
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


@app.post("/spam-reports")
@traceable(name="report_create")
def create_report(req: ReportReq, user=Depends(get_current_user)):
    text = req.email_content.strip()
    if not text:
        raise HTTPException(status_code=400, detail="?댁슜???낅젰?댁＜?몄슂")
    print(f"[LangSmith][path] /spam-reports received username={user['username']} text={text[:80]!r}")
    log_report_input_to_langsmith(text=text, username=user["username"])
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT value FROM system_settings WHERE key='report_enabled'")
    row = cur.fetchone()
    if row and row["value"] == "false":
        conn.close()
        raise HTTPException(status_code=403, detail="?꾩옱 ?곷떞 ?붿껌??鍮꾪솢?깊솕?섏뼱 ?덉뒿?덈떎")
    cur = conn.cursor()
    cur.execute(
        "INSERT INTO spam_reports (user_id, email_content) VALUES (?, ?)",
        (user["id"], text)
    )
    conn.commit()
    conn.close()
    try:
        from langchain_core.tracers.langchain import wait_for_all_tracers
        wait_for_all_tracers()
        print("[LangSmith][flush] wait_for_all_tracers completed for /spam-reports")
    except Exception as e:
        print(f"[LangSmith][flush] helper unavailable for /spam-reports: {e}")
    return {"message": "?좉퀬 ?묒닔 ?꾨즺"}


@app.patch("/spam-reports/{report_id}")
def update_report(report_id: int, req: ReportUpdateReq, user=Depends(require_role("admin", "counselor"))):
    conn = get_db()
    cur = conn.cursor()
    # 泥섎━ ?꾩뿉 ?대찓???댁슜 議고쉶 (vectorstore ?낅뜲?댄듃??
    cur.execute("SELECT email_content, status FROM spam_reports WHERE id=?", (report_id,))
    report = cur.fetchone()

    cur.execute(
        "UPDATE spam_reports SET status=?, counselor_note=?, counselor_id=? WHERE id=?",
        (req.status, req.counselor_note, user["id"], report_id)
    )
    # ?꾨즺 泥섎━ ???ㅼ썙??異붽?
    if req.keywords:
        for kw in req.keywords:
            cur.execute("INSERT OR IGNORE INTO spam_keywords (keyword) VALUES (?)", (kw,))
    conn.commit()
    conn.close()

    # ?꾨즺 泥섎━ ??RAG vectorstore???숈뒿 ?덉떆 異붽? (吏???숈뒿)
    if report and req.status in ("spam", "ham", "completed") and _rag_available:
        label = "spam" if req.status == "spam" else "ham"
        # status媛 completed硫?counselor_note?먯꽌 ?먯젙 ?좎텛 (?놁쑝硫?ham?쇰줈)
        if req.status == "completed":
            note = (req.counselor_note or "").lower()
            label = "spam" if "?ㅽ뙵" in note else "ham"
        if ALLOW_REPORT_VECTORSTORE_LEARN:
            try:
                added = add_to_vectorstore(report["email_content"], label)
                if added:
                    print(f"??vectorstore ?낅뜲?댄듃: [{label}] report_id={report_id}")
            except Exception as e:
                print(f"?좑툘  vectorstore ?낅뜲?댄듃 ?ㅽ뙣: {e}")
        else:
            print(f"[SAFEGUARD] report-based vectorstore learning blocked for report_id={report_id}")

    return {"message": "?낅뜲?댄듃 ?꾨즺"}


# ?????????????????????????????
# SYSTEM SETTINGS
# ?????????????????????????????
@app.get("/settings")
def get_settings(user=Depends(get_current_user)):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT key, value FROM system_settings")
    rows = cur.fetchall()
    conn.close()
    return {r["key"]: r["value"] for r in rows}


class SettingUpdateReq(BaseModel):
    value: str


@app.patch("/settings/{key}")
def update_setting(key: str, req: SettingUpdateReq, user=Depends(require_role("admin"))):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("INSERT OR REPLACE INTO system_settings (key, value) VALUES (?, ?)", (key, req.value))
    conn.commit()
    conn.close()
    return {"message": "?ㅼ젙 蹂寃??꾨즺"}


# ?????????????????????????????
# SPAM KEYWORDS
# ?????????????????????????????
@app.get("/spam-keywords")
def get_keywords(user=Depends(get_current_user)):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT * FROM spam_keywords ORDER BY created_at DESC")
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


class KeywordAddReq(BaseModel):
    keyword: str


@app.post("/spam-keywords")
def add_keyword(req: KeywordAddReq, user=Depends(require_role("admin", "counselor"))):
    kw = req.keyword.strip()
    if not kw:
        raise HTTPException(status_code=400, detail="?ㅼ썙?쒕? ?낅젰?댁＜?몄슂")
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT id FROM spam_keywords WHERE keyword=?", (kw,))
    if cur.fetchone():
        conn.close()
        raise HTTPException(status_code=400, detail="?대? ?깅줉???ㅼ썙?쒖엯?덈떎")
    cur.execute("INSERT INTO spam_keywords (keyword) VALUES (?)", (kw,))
    conn.commit()
    new_id = cur.lastrowid
    conn.close()
    return {"message": "?깅줉 ?꾨즺", "id": new_id, "keyword": kw}


@app.delete("/spam-keywords/{kw_id}")
def delete_keyword(kw_id: int, user=Depends(require_role("admin", "counselor"))):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("DELETE FROM spam_keywords WHERE id=?", (kw_id,))
    conn.commit()
    conn.close()
    return {"message": "??젣 ?꾨즺"}


# ?????????????????????????????
# CHAT
# ?????????????????????????????
class ChatReq(BaseModel):
    message: str


def _build_model_reply(model_name: str, result: dict) -> dict:
    if not result:
        return {
            "model": model_name,
            "reply": f"[{model_name}] 분석 결과를 받지 못했습니다.",
            "is_spam": False,
            "error": True,
        }

    if result.get("error"):
        return {
            "model": model_name,
            "reply": f"[{model_name}] ?먮퀎 ?ㅽ뙣: {result['error']}",
            "is_spam": False,
            "error": True,
        }

    verdict = "스팸" if result.get("is_spam") else "정상"
    raw_answer = (result.get("raw_answer") or "").strip()
    if raw_answer:
        reply = f"[{model_name}] {verdict}으로 판단했습니다.\n{raw_answer}"
    else:
        reason = "스팸 징후가 감지되었습니다." if result.get("is_spam") else "정상 메시지로 보입니다."
        reply = f"[{model_name}] {verdict}으로 판단했습니다. {reason}"

    return {
        "model": model_name,
        "reply": reply,
        "is_spam": bool(result.get("is_spam")),
        "error": False,
    }


@traceable(name="chat_classify_with_both_models")
def _classify_with_both_models(text: str) -> list[dict]:
    print(f"[LangSmith][path] /chat -> _classify_with_both_models text_len={len(text)}")
    if not _rag_available:
        is_spam = detect_spam(text)
        fallback_verdict = "스팸" if is_spam else "정상"
        return [
            {
                "model": "GPT",
                "reply": f"[GPT] {fallback_verdict}으로 판단했습니다. 현재 GPT 엔진을 사용할 수 없어 규칙 기반 결과를 대신 보여줍니다.",
                "is_spam": is_spam,
                "error": True,
            },
            {
                "model": "Qwen",
                "reply": f"[Qwen] {fallback_verdict}으로 판단했습니다. 현재 Qwen 엔진을 사용할 수 없어 규칙 기반 결과를 대신 보여줍니다.",
                "is_spam": is_spam,
                "error": True,
            },
        ]

    from concurrent.futures import ThreadPoolExecutor
    from rag_engine import _retrieve

    examples = _retrieve(text)
    with ThreadPoolExecutor(max_workers=2) as executor:
        gpt_future = executor.submit(classify_with_gpt, text, examples)
        qwen_future = executor.submit(classify_with_qwen, text, examples)
        gpt_result = gpt_future.result()
        qwen_result = qwen_future.result()

    return [
        _build_model_reply("GPT", gpt_result),
        _build_model_reply("Qwen", qwen_result),
    ]


SPAM_LIST_KEYWORDS = ["스팸 단어", "스팸 목록", "스팸 키워드", "spam list", "스팸 리스트", "등록된 키워드", "키워드 목록"]

@app.post("/chat")
@traceable(name="chat_endpoint")
def chat(req: ChatReq, user=Depends(get_current_user)):
    text = req.message.strip()
    if not text:
        raise HTTPException(status_code=400, detail="메시지를 입력해주세요")
    print(f"[LangSmith][path] /chat received username={user['username']} text={text[:80]!r}")
    log_chat_input_to_langsmith(text=text, username=user["username"])

    conn = get_db()
    cur = conn.cursor()

    text_lower = text.lower()
    if any(kw in text_lower for kw in SPAM_LIST_KEYWORDS):
        cur.execute("SELECT keyword FROM spam_keywords ORDER BY id DESC LIMIT 50")
        kws = [row["keyword"] for row in cur.fetchall()]
        if kws:
            reply = f"등록된 스팸 키워드 {len(kws)}개:\n" + ", ".join(kws)
        else:
            reply = "등록된 스팸 키워드가 없습니다."
        cur.execute(
            "INSERT INTO chat_logs (user_id, username, role, message, is_spam) VALUES (?,?,?,?,?)",
            (user["id"], user["username"], "user", text, 0),
        )
        cur.execute(
            "INSERT INTO chat_logs (user_id, username, role, message, is_spam) VALUES (?,?,?,?,?)",
            (user["id"], user["username"], "ai", reply, 0),
        )
        conn.commit()
        conn.close()
        return {"reply": reply, "is_spam": False, "replies": [{"model": "System", "reply": reply, "is_spam": False}]}

    replies = _classify_with_both_models(text)
    final_is_spam = any(item["is_spam"] for item in replies if not item.get("error"))
    if not any(not item.get("error") for item in replies):
        final_is_spam = any(item["is_spam"] for item in replies)

    cur.execute(
        "INSERT INTO chat_logs (user_id, username, role, message, is_spam) VALUES (?,?,?,?,?)",
        (user["id"], user["username"], "user", text, 1 if final_is_spam else 0),
    )
    for item in replies:
        cur.execute(
            "INSERT INTO chat_logs (user_id, username, role, message, is_spam) VALUES (?,?,?,?,?)",
            (user["id"], user["username"], "ai", item["reply"], 1 if item["is_spam"] else 0),
        )

    conn.commit()
    conn.close()
    try:
        from langchain_core.tracers.langchain import wait_for_all_tracers
        wait_for_all_tracers()
        print("[LangSmith][flush] wait_for_all_tracers completed for /chat")
    except Exception as e:
        print(f"[LangSmith][flush] helper unavailable: {e}")
    return {"reply": replies[0]["reply"], "is_spam": final_is_spam, "replies": replies}


@app.post("/chat/file")
async def chat_file(file: UploadFile = File(...), user=Depends(get_current_user)):
    content = await file.read()
    filename = file.filename or "파일"
    text = ""
    extracted_via = "unknown"

    fn = filename.lower()
    try:
        if fn.endswith(".pdf"):
            import pdfplumber
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    text += (page.extract_text() or "")
            extracted_via = "pdfplumber"
            if not text.strip():
                try:
                    import pypdf
                    reader = pypdf.PdfReader(io.BytesIO(content))
                    text = "\n".join((page.extract_text() or "") for page in reader.pages)
                    extracted_via = "pypdf"
                except Exception:
                    pass
        elif fn.endswith(".docx"):
            import docx
            doc = docx.Document(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs)
            extracted_via = "docx"
        elif fn.endswith(".xlsx") or fn.endswith(".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    text += " ".join(str(c) for c in row if c is not None) + "\n"
            extracted_via = "openpyxl"
        elif fn.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            extracted_via = "pptx"
        elif fn.endswith(".hwpx"):
            import zipfile, xml.etree.ElementTree as ET
            with zipfile.ZipFile(io.BytesIO(content)) as z:
                for name in z.namelist():
                    if name.startswith("Contents/") and name.endswith(".xml"):
                        root = ET.fromstring(z.read(name))
                        for elem in root.iter():
                            if elem.text:
                                text += elem.text + " "
            extracted_via = "hwpx"
        elif fn.endswith(".hwp"):
            import subprocess, tempfile
            with tempfile.NamedTemporaryFile(suffix=".hwp", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                backend_dir = os.path.dirname(os.path.abspath(__file__))
                env = os.environ.copy()
                env["PYTHONPATH"] = backend_dir + os.pathsep + env.get("PYTHONPATH", "")
                result = subprocess.run(
                    [sys.executable, "-m", "hwp5.hwp5txt", tmp_path],
                    capture_output=True,
                    timeout=30,
                    cwd=backend_dir,
                    env=env,
                )
                text = _decode_bytes_best_effort(result.stdout).strip()
                extracted_via = "hwp5txt"
                if not text:
                    preview_text = _extract_text_from_hwp_preview(content)
                    if preview_text:
                        text = preview_text
                        extracted_via = "hwp-preview"
                if not text:
                    body_text = _extract_text_from_hwp_body(content)
                    if body_text:
                        text = body_text
                        extracted_via = "hwp-body"
                if not text and result.stderr:
                    raise Exception(_decode_bytes_best_effort(result.stderr)[:200])
            finally:
                os.unlink(tmp_path)
        elif fn.endswith(".eml"):
            text = _extract_text_from_eml(content)
            extracted_via = "eml-parser"
        else:
            text = _decode_bytes_best_effort(content)
            extracted_via = "best-effort-bytes"
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"파일 읽기 실패: {str(e)}")

    text = _normalize_extracted_text(text)
    print(
        f"[file-upload] filename={filename!r} extracted_len={len(text)} "
        f"hangul_ratio={_hangul_ratio(text):.3f} via={extracted_via} preview={text[:80]!r}"
    )
    if not text:
        ext = os.path.splitext(filename)[1].lower() or "unknown"
        if ext in {".pdf", ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}:
            detail = f"파일에서 텍스트를 추출할 수 없습니다. 스캔형 {ext} 파일일 가능성이 큽니다."
        else:
            detail = f"파일에서 텍스트를 추출할 수 없습니다. 확장자={ext}, 추출기={extracted_via}"
        raise HTTPException(status_code=400, detail=detail)

    replies = _classify_with_both_models(text)
    final_is_spam = any(item["is_spam"] for item in replies if not item.get("error"))
    if not any(not item.get("error") for item in replies):
        final_is_spam = any(item["is_spam"] for item in replies)

    conn = get_db()
    cur = conn.cursor()
    cur.execute(
        "INSERT INTO chat_logs (user_id, username, role, message, is_spam) VALUES (?,?,?,?,?)",
        (user["id"], user["username"], "user", f"[파일업로드] {filename}", 1 if final_is_spam else 0),
    )
    for item in replies:
        body = item["reply"].split("] ", 1)[1] if "] " in item["reply"] else item["reply"]
        file_reply = f"[{item['model']}] [{filename}] {body}"
        cur.execute(
            "INSERT INTO chat_logs (user_id, username, role, message, is_spam) VALUES (?,?,?,?,?)",
            (user["id"], user["username"], "ai", file_reply, 1 if item["is_spam"] else 0),
        )
        item["reply"] = file_reply

    conn.commit()
    conn.close()
    return {"reply": replies[0]["reply"], "is_spam": final_is_spam, "replies": replies, "filename": filename, "text_preview": text[:500]}

@app.get("/chat/history")
def chat_history(user=Depends(get_current_user)):
    conn = get_db()
    cur = conn.cursor()
    cur.execute(
        "SELECT * FROM chat_logs WHERE user_id=? ORDER BY created_at ASC",
        (user["id"],)
    )
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


@app.get("/chat/all")
def chat_all(user=Depends(require_role("admin", "counselor"))):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT * FROM chat_logs ORDER BY created_at DESC")
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


# ?????????????????????????????
# COUNSELOR REVIEWS
# ?????????????????????????????
class ReviewReq(BaseModel):
    report_id: int
    stars: int
    comment: Optional[str] = ""
    accuracy_stars: Optional[int] = None
    processing_stars: Optional[int] = None
    clarity_stars: Optional[int] = None
    speed_stars: Optional[int] = None


@app.post("/counselor-reviews")
def create_review(req: ReviewReq, user=Depends(get_current_user)):
    if not (1 <= req.stars <= 5):
        raise HTTPException(status_code=400, detail="별점은 1~5 사이여야 합니다")
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT counselor_id FROM spam_reports WHERE id=? AND user_id=?", (req.report_id, user["id"]))
    row = cur.fetchone()
    if not row:
        conn.close()
        raise HTTPException(status_code=404, detail="?좉퀬瑜?李얠쓣 ???놁뒿?덈떎")
    cur.execute(
        "INSERT OR IGNORE INTO counselor_reviews (report_id, reviewer_id, counselor_id, stars, comment, accuracy_stars, processing_stars, clarity_stars, speed_stars) VALUES (?,?,?,?,?,?,?,?,?)",
        (req.report_id, user["id"], row["counselor_id"], req.stars, req.comment,
         req.accuracy_stars, req.processing_stars, req.clarity_stars, req.speed_stars)
    )
    conn.commit()
    conn.close()
    return {"message": "?됯? ?꾨즺"}


@app.get("/counselor-reviews")
def get_reviews(user=Depends(require_role("admin"))):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("""
        SELECT rv.*, u1.username as reviewer, u2.username as counselor_name, r.email_content
        FROM counselor_reviews rv
        JOIN users u1 ON rv.reviewer_id = u1.id
        LEFT JOIN users u2 ON rv.counselor_id = u2.id
        JOIN spam_reports r ON rv.report_id = r.id
        ORDER BY rv.created_at DESC
    """)
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


@app.get("/counselor-reviews/mine")
def get_my_reviews(user=Depends(require_role("counselor"))):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("""
        SELECT rv.*, u.username as reviewer, r.email_content
        FROM counselor_reviews rv
        JOIN users u ON rv.reviewer_id = u.id
        JOIN spam_reports r ON rv.report_id = r.id
        WHERE rv.counselor_id = ?
        ORDER BY rv.created_at DESC
    """, (user["id"],))
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


# ?????????????????????????????
# USER SPAM RULES
# ?????????????????????????????
class SpamRuleReq(BaseModel):
    rule_type: str
    value: str


@app.get("/user-spam-rules")
def get_spam_rules(user=Depends(get_current_user)):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT * FROM user_spam_rules WHERE user_id=? ORDER BY created_at DESC", (user["id"],))
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


@app.post("/user-spam-rules")
def add_spam_rule(req: SpamRuleReq, user=Depends(get_current_user)):
    if req.rule_type not in ("keyword", "sentence", "email"):
        raise HTTPException(status_code=400, detail="?섎せ??洹쒖튃 ??낆엯?덈떎")
    if not req.value.strip():
        raise HTTPException(status_code=400, detail="媛믪쓣 ?낅젰?댁＜?몄슂")
    conn = get_db()
    cur = conn.cursor()
    cur.execute(
        "INSERT INTO user_spam_rules (user_id, rule_type, value) VALUES (?,?,?)",
        (user["id"], req.rule_type, req.value.strip())
    )
    conn.commit()
    conn.close()
    return {"message": "洹쒖튃 異붽? ?꾨즺"}


@app.delete("/user-spam-rules/{rule_id}")
def delete_spam_rule(rule_id: int, user=Depends(get_current_user)):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("DELETE FROM user_spam_rules WHERE id=? AND user_id=?", (rule_id, user["id"]))
    conn.commit()
    conn.close()
    return {"message": "??젣 ?꾨즺"}


# ?????????????????????????????
# AI ?깅뒫 ?됯?
# ?????????????????????????????
class EvalItem(BaseModel):
    text: str
    label: str  # 'spam' or 'ham'

def _run_evaluation(items):
    tp = tn = fp = fn = 0
    results = []
    for item in items:
        predicted = detect_spam(item["text"])
        actual = item["label"] == "spam"
        if actual and predicted:          tp += 1
        elif not actual and not predicted: tn += 1
        elif not actual and predicted:     fp += 1
        elif actual and not predicted:     fn += 1
        results.append({
            "text": item["text"][:120],
            "label": item["label"],
            "predicted": "spam" if predicted else "ham",
            "correct": predicted == actual,
        })
    total = len(items)
    accuracy  = (tp + tn) / total if total else 0
    precision = tp / (tp + fp) if (tp + fp) else 0
    recall    = tp / (tp + fn) if (tp + fn) else 0
    f1        = 2 * precision * recall / (precision + recall) if (precision + recall) else 0
    return {
        "total": total,
        "accuracy": round(accuracy, 4),
        "precision": round(precision, 4),
        "recall": round(recall, 4),
        "f1": round(f1, 4),
        "tp": tp, "tn": tn, "fp": fp, "fn": fn,
        "results": results,
    }

@app.post("/ai/evaluate")
def ai_evaluate(items: list[EvalItem], user=Depends(require_role("admin"))):
    if not items:
        raise HTTPException(status_code=400, detail="?됯? ?곗씠?곌? ?놁뒿?덈떎")
    data = [{"text": it.text, "label": it.label.strip().lower()} for it in items]
    return _run_evaluation(data)

@app.get("/ai/evaluate/auto")
def ai_evaluate_auto(user=Depends(require_role("admin"))):
    conn = get_db()
    cur = conn.cursor()
    cur.execute(
        "SELECT message, is_spam FROM chat_logs WHERE role='user' AND is_spam IS NOT NULL ORDER BY created_at DESC LIMIT 500"
    )
    rows = cur.fetchall()
    conn.close()
    if not rows:
        raise HTTPException(status_code=404, detail="?됯????덉씠釉붾맂 ?곗씠?곌? ?놁뒿?덈떎.")
    data = [{"text": r["message"], "label": "spam" if r["is_spam"] == 1 else "ham"} for r in rows]
    result = _run_evaluation(data)
    result["source"] = "chat_logs"
    result["note"] = f"chat_logs에서 수집한 {len(data)}개의 레이블 데이터를 평가했습니다."
    return result

@app.get("/ai/evaluate/testset")
def ai_evaluate_testset(split: str = "test", user=Depends(require_role("admin"))):
    if split not in ("test", "val"):
        raise HTTPException(status_code=400, detail="split은 'test' 또는 'val' 이어야 합니다")
    conn = get_db()
    cur = conn.cursor()
    # eval_testset ?뚯씠釉??놁쑝硫??덈궡
    try:
        cur.execute("SELECT text, label FROM eval_testset WHERE split=?", (split,))
        rows = cur.fetchall()
    except Exception:
        conn.close()
        raise HTTPException(status_code=404, detail="?뚯뒪???곗씠?곌? ?놁뒿?덈떎. split_data.py瑜?癒쇱? ?ㅽ뻾?댁＜?몄슂.")
    conn.close()
    if not rows:
        raise HTTPException(status_code=404, detail=f"{split} ?곗씠?곌? ?놁뒿?덈떎. split_data.py瑜?癒쇱? ?ㅽ뻾?댁＜?몄슂.")
    data = [{"text": r["text"], "label": r["label"]} for r in rows]
    result = _run_evaluation(data)
    result["source"] = split
    result["note"] = f"{'테스트' if split == 'test' else '검증'} 데이터셋 {len(data)}건을 평가했습니다."
    return result


# ?????????????????????????????
# RAG 吏곸젒 遺꾨쪟 ?붾뱶?ъ씤??(愿由ъ옄??
# ?????????????????????????????
class RagRequest(BaseModel):
    text: str
    model: str = "both"  # 湲곕낯媛? gpt + qwen ?숈떆 ?ㅽ뻾

@app.post("/rag/fast-classify")
def rag_fast_classify(req: RagRequest, user=Depends(require_role("admin"))):
    """
    ?꾨쿋???ы몴留??ъ슜??利됯컖 ?먯젙 (~0.5珥? GPT/Qwen 誘몄궗??
    """
    if not _rag_available:
        raise HTTPException(status_code=503, detail="RAG ?붿쭊???ъ슜?????놁뒿?덈떎.")
    if not req.text.strip():
        raise HTTPException(status_code=400, detail="?띿뒪?몃? ?낅젰?댁＜?몄슂.")
    from rag_engine import fast_classify, _retrieve
    examples = _retrieve(req.text, k=7)
    result = fast_classify(req.text, k=7)
    spam_votes = sum(1 for e in examples if e["label"] == "spam")
    ham_votes = len(examples) - spam_votes
    return {
        "is_spam": bool(result) if result is not None else None,
        "method": "embedding_vote",
        "spam_votes": spam_votes,
        "ham_votes": ham_votes,
        "total_votes": len(examples),
        "confidence": round(max(spam_votes, ham_votes) / len(examples), 2) if examples else 0,
        "retrieved_examples": [
            {"text": e["text"][:80], "label": e["label"], "distance": round(e["distance"], 3)}
            for e in examples[:5]
        ]
    }


@app.post("/rag/classify")
def rag_classify(req: RagRequest, user=Depends(require_role("admin", "counselor"))):
    if not _rag_available:
        raise HTTPException(status_code=503, detail="RAG ?붿쭊???ъ슜?????놁뒿?덈떎.")
    if not req.text.strip():
        raise HTTPException(status_code=400, detail="?띿뒪?몃? ?낅젰?댁＜?몄슂.")

    if req.model == "gpt":
        return classify_with_gpt(req.text)
    elif req.model == "qwen":
        return classify_with_qwen(req.text)
    else:  # "both" - ?꾨쿋??1??怨듭쑀 ??GPT + Qwen ?숈떆 ?ㅽ뻾
        from rag_engine import _retrieve
        from concurrent.futures import ThreadPoolExecutor
        examples = _retrieve(req.text)  # ?꾨쿋??1踰덈쭔 怨꾩궛
        with ThreadPoolExecutor(max_workers=2) as ex:
            f_gpt  = ex.submit(classify_with_gpt,  req.text, examples)
            f_qwen = ex.submit(classify_with_qwen, req.text, examples)
            gpt_result  = f_gpt.result()
            qwen_result = f_qwen.result()
        both_spam = gpt_result.get("is_spam", False) or qwen_result.get("is_spam", False)
        return {
            "gpt": gpt_result,
            "qwen": qwen_result,
            "final_is_spam": both_spam,
            "agreement": gpt_result.get("is_spam") == qwen_result.get("is_spam")
        }


class RagLearnReq(BaseModel):
    text: str
    label: str  # 'spam' or 'ham'

@app.post("/rag/learn")
def rag_learn(req: RagLearnReq, user=Depends(require_role("admin", "counselor"))):
    """?섎룞?쇰줈 ?덉떆瑜?vectorstore??異붽?"""
    _ensure_learning_enabled(
        ALLOW_RAG_DIRECT_LEARN,
        "직접 vectorstore 학습은 기본 비활성화되어 있습니다. 누수 방지를 위해 명시적으로 허용해야 합니다.",
    )
    if req.label not in ("spam", "ham"):
        raise HTTPException(status_code=400, detail="label? 'spam' ?먮뒗 'ham'?댁뼱???⑸땲??")
    if not _rag_available:
        raise HTTPException(status_code=503, detail="RAG ?붿쭊???ъ슜?????놁뒿?덈떎.")
    added = add_to_vectorstore(req.text, req.label)
    if not added:
        raise HTTPException(status_code=500, detail="vectorstore 異붽? ?ㅽ뙣")
    return {"message": f"[{req.label}] ?덉떆 異붽? ?꾨즺"}


@app.post("/rag/train")
def rag_train(user=Depends(require_role("admin"))):
    """
    vectorstore ?щ퉴??
    train_split.csv + training_corrections(DB ?꾩쟻 媛쒖꽑 ?덉떆) ?⑹퀜???숈뒿
    ?꾨쿋??罹먯떆 ?ъ슜 ??train_split? ?ы샇異??놁쓬
    """
    try:
        import pandas as pd
        import faiss
        import numpy as np
        import pickle
        from openai import OpenAI
        from dotenv import load_dotenv

        BASE = os.path.dirname(__file__)
        load_dotenv(os.path.join(BASE, ".env"))
        oai = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        VS_DIR = os.path.join(BASE, "vectorstore")
        CACHE_EMB = os.path.join(VS_DIR, "embeddings_cache.npy")
        CACHE_META = os.path.join(VS_DIR, "cache_meta.pkl")
        os.makedirs(VS_DIR, exist_ok=True)

        train_csv = os.path.join(BASE, "..", "train_split.csv")
        if not os.path.exists(train_csv):
            raise HTTPException(status_code=400, detail="train_split.csv가 없습니다. split_data.py를 먼저 실행해주세요.")
        df = pd.read_csv(train_csv).dropna(subset=["text", "label"])
        texts = df["text"].tolist()
        labels = df["label"].tolist()
        print(f"학습 데이터 {len(texts)}건")

        # 임베딩 캐시 확인
        if os.path.exists(CACHE_EMB) and os.path.exists(CACHE_META):
            with open(CACHE_META, "rb") as f:
                cached = pickle.load(f)
            if cached.get("texts") == texts:
                print("캐시된 임베딩을 재사용합니다.")
                embeddings = np.load(CACHE_EMB)
            else:
                print("임베딩을 새로 생성하는 중입니다...")
                embeddings = _get_embeddings(oai, texts)
                np.save(CACHE_EMB, embeddings)
                with open(CACHE_META, "wb") as f:
                    pickle.dump({"texts": texts}, f)
        else:
            print("임베딩을 생성하는 중입니다...")
            embeddings = _get_embeddings(oai, texts)
            np.save(CACHE_EMB, embeddings)
            with open(CACHE_META, "wb") as f:
                pickle.dump({"texts": texts}, f)

        # FAISS 인덱스 빌드
        dim = embeddings.shape[1]
        index = faiss.IndexFlatL2(dim)
        index.add(embeddings)

        index_bytes = faiss.serialize_index(index)
        with open(os.path.join(VS_DIR, "spam.index"), "wb") as f:
            f.write(index_bytes)
        with open(os.path.join(VS_DIR, "metadata.pkl"), "wb") as f:
            pickle.dump({"texts": texts, "labels": labels}, f)

        # rag_engine 리로드
        import rag_engine
        rag_engine._load_vectorstore()

        # training_corrections 중 train 학습에서 나온 보정만 다시 반영
        conn = get_db()
        cur = conn.cursor()
        try:
            cur.execute("SELECT text, label FROM training_corrections WHERE source='improve_train'")
            corrections = cur.fetchall()
        except Exception:
            corrections = []
        conn.close()

        added = 0
        for row in corrections:
            if add_to_vectorstore(row["text"], row["label"]):
                added += 1

        total = rag_engine._index.ntotal if rag_engine._index else 0
        print(f"?숈뒿 ?꾨즺: 踰≫꽣 {total}媛?(媛쒖꽑 ?덉떆 {added}媛??ы븿)")
        return {"message": "?숈뒿 ?꾨즺", "vector_count": total, "corrections_added": added}

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


def _get_embeddings(client, texts: list, batch_size=100):
    import numpy as np
    all_emb = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i+batch_size]
        resp = client.embeddings.create(model="text-embedding-3-small", input=batch)
        all_emb.extend([e.embedding for e in resp.data])
        print(f"  ?꾨쿋?? {min(i+batch_size, len(texts))}/{len(texts)}")
    return np.array(all_emb, dtype="float32")


def _eval_with_model(rows, model_fn, model_name: str, eval_debug: dict):
    _log_eval_dataset(
        mode=eval_debug["mode"],
        split=eval_debug["split"],
        eval_csv_path=eval_debug["eval_csv_path"],
        row_count=eval_debug["row_count"],
        first_3_titles=eval_debug["first_3_titles"],
        effective_row_count=eval_debug["effective_row_count"],
        model_path=model_name,
    )
    tp = tn = fp = fn = 0
    fp_list = []  # 오탐: 정상인데 스팸으로 잘못 판정
    fn_list = []  # 미탐: 스팸인데 정상으로 놓침
    for row in rows:
        actual_spam = (row["label"] == "spam")
        try:
            result = model_fn(row["text"])
            predicted_spam = result.get("is_spam", False)
        except Exception:
            predicted_spam = False
        text_preview = row["text"][:120]
        if actual_spam and predicted_spam:
            tp += 1
        elif not actual_spam and not predicted_spam:
            tn += 1
        elif not actual_spam and predicted_spam:
            fp += 1
            fp_list.append(text_preview)
        elif actual_spam and not predicted_spam:
            fn += 1
            fn_list.append(text_preview)
    total = tp + tn + fp + fn
    accuracy  = (tp + tn) / total if total else 0
    precision = tp / (tp + fp) if (tp + fp) else 0
    recall    = tp / (tp + fn) if (tp + fn) else 0
    f1        = 2 * precision * recall / (precision + recall) if (precision + recall) else 0
    return {
        "total": total, "tp": tp, "tn": tn, "fp": fp, "fn": fn,
        "accuracy": round(accuracy, 4), "precision": round(precision, 4),
        "recall": round(recall, 4), "f1": round(f1, 4),
        "fp_list": fp_list,
        "fn_list": fn_list,
    }

@app.get("/rag/val")
def rag_val(split: str = "val", max_samples: int = 120, mode: str = "fast", user=Depends(require_role("admin"))):
    """
    val 또는 test 데이터셋으로 평가
    split: 'train' | 'val' | 'test'
    mode: 'fast' (임베딩 기반, 빠름) | 'llm' (GPT+Qwen, 느림)
    """
    if not _rag_available:
        raise HTTPException(status_code=503, detail="RAG 서비스를 사용할 수 없습니다.")

    eval_data = _load_eval_dataset(split, max_samples)
    rows = eval_data["rows"]

    if not rows:
        raise HTTPException(status_code=404, detail=f"{split} 데이터가 없습니다.")

    if mode == "fast":
        from rag_engine import fast_classify
        _log_eval_dataset(
            mode=mode,
            split=split,
            eval_csv_path=eval_data["eval_csv_path"],
            row_count=eval_data["row_count"],
            first_3_titles=eval_data["first_3_titles"],
            effective_row_count=eval_data["effective_row_count"],
            model_path="fast_classify",
        )
        tp = tn = fp = fn = 0
        fp_list = []
        fn_list = []
        with _evaluation_write_blocked():
            for row in rows:
                actual_spam = (row["label"] == "spam")
                text = row["text"]
                result = fast_classify(text)
                predicted_spam = bool(result) if result is not None else False
                text_preview = text[:120]
                if actual_spam and predicted_spam:
                    tp += 1
                elif not actual_spam and not predicted_spam:
                    tn += 1
                elif not actual_spam and predicted_spam:
                    fp += 1
                    fp_list.append(text_preview)
                elif actual_spam and not predicted_spam:
                    fn += 1
                    fn_list.append(text_preview)
        total = tp + tn + fp + fn
        accuracy  = (tp + tn) / total if total else 0
        precision = tp / (tp + fp) if (tp + fp) else 0
        recall    = tp / (tp + fn) if (tp + fn) else 0
        f1        = 2 * precision * recall / (precision + recall) if (precision + recall) else 0
        fast_result = {
            "model": "FastEmbed",
            "total": total, "tp": tp, "tn": tn, "fp": fp, "fn": fn,
            "accuracy": round(accuracy, 4), "precision": round(precision, 4),
            "recall": round(recall, 4), "f1": round(f1, 4),
            "fp_list": fp_list,
            "fn_list": fn_list,
        }
        return {
            "split": split,
            "mode": "fast",
            "gpt": fast_result,
            "qwen": fast_result,
            "eval_debug": {
                **eval_data,
                "mode": mode,
            },
        }

    eval_debug = {
        **eval_data,
        "mode": mode,
        "split": split,
    }
    with _evaluation_write_blocked():
        gpt_result  = _eval_with_model(rows, classify_with_gpt, "classify_with_gpt", eval_debug)
        qwen_result = _eval_with_model(rows, classify_with_qwen, "classify_with_qwen", eval_debug)

    return {
        "split": split,
        "mode": "llm",
        "gpt":  {"model": "GPT-4o-mini",  **gpt_result},
        "qwen": {"model": "Qwen2.5-7b",   **qwen_result},
        "eval_debug": eval_debug,
    }


@app.post("/rag/improve")
def rag_improve(
    split: str = "val",
    mode: str = "fast",
    max_samples: int = 500,
    dry_run: bool = False,
    user=Depends(require_role("admin"))
):
    """
    媛쒖꽑 猷⑦봽 1???ㅽ뻾:
    1. val/test set ?됯?
    2. ?ㅻ텇瑜??덉떆 ??vectorstore???뺣떟 ?덉씠釉붾줈 異붽?
    3. 媛쒖꽑 寃곌낵 諛섑솚 + DB ?대젰 ???
    mode: 'fast' (?꾨쿋???ы몴, 臾대즺) | 'gpt' (GPT ?몄텧, ?좊즺)
    """
    if not _rag_available:
        raise HTTPException(status_code=503, detail="RAG ?붿쭊???ъ슜?????놁뒿?덈떎.")

    eval_data = _load_eval_dataset(split, max_samples)
    rows = eval_data["rows"]

    if not rows:
        raise HTTPException(status_code=404, detail=f"{split} ?곗씠???놁쓬")

    from rag_engine import fast_classify
    _log_eval_dataset(
        mode=mode,
        split=split,
        eval_csv_path=eval_data["eval_csv_path"],
        row_count=eval_data["row_count"],
        first_3_titles=eval_data["first_3_titles"],
        effective_row_count=eval_data["effective_row_count"],
        model_path="fast_classify" if mode == "fast" else "classify_with_gpt",
    )

    tp = tn = fp = fn = 0
    misclassified = []

    with _evaluation_write_blocked():
        for row in rows:
            text, true_label = row["text"], row["label"]
            actual_spam = (true_label == "spam")

            if mode == "fast":
                result = fast_classify(text)
                if result is None:
                    continue
                predicted_spam = result
            else:
                try:
                    result = classify_with_gpt(text)
                    predicted_spam = result.get("is_spam", False)
                except Exception:
                    predicted_spam = False

            if actual_spam and predicted_spam:       tp += 1
            elif not actual_spam and not predicted_spam: tn += 1
            elif not actual_spam and predicted_spam:     fp += 1; misclassified.append((text, "ham"))
            elif actual_spam and not predicted_spam:     fn += 1; misclassified.append((text, "spam"))

    total = tp + tn + fp + fn
    accuracy  = (tp + tn) / total if total else 0
    precision = tp / (tp + fp) if (tp + fp) else 0
    recall    = tp / (tp + fn) if (tp + fn) else 0
    f1        = 2 * precision * recall / (precision + recall) if (precision + recall) else 0

    # 검증/테스트 데이터는 평가 전용으로만 사용한다.
    added = 0
    improve_note = ""
    learned = False
    if not dry_run and split == "train":
        learned = True
        conn2 = get_db()
        cur2 = conn2.cursor()
        for text, correct_label in misclassified:
            if add_to_vectorstore(text, correct_label):
                added += 1
            cur2.execute(
                "INSERT INTO training_corrections (text, label, source) VALUES (?, ?, 'improve_train')",
                (text, correct_label)
            )
        conn2.commit()
        conn2.close()
    elif split in ("val", "test"):
        improve_note = f"{split} split은 평가 전용이라 vectorstore에 반영하지 않았습니다."
    elif dry_run:
        improve_note = "dry_run=True라 vectorstore에 반영하지 않았습니다."

    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT COALESCE(MAX(iteration), 0) + 1 FROM improvement_history")
    next_iter = cur.fetchone()[0]
    cur.execute("""
        INSERT INTO improvement_history
        (iteration, mode, accuracy, precision_score, recall, f1, total, tp, tn, fp, fn, added_count)
        VALUES (?,?,?,?,?,?,?,?,?,?,?,?)
    """, (next_iter, mode, round(accuracy,4), round(precision,4), round(recall,4), round(f1,4),
          total, tp, tn, fp, fn, added))
    conn.commit()
    conn.close()

    return {
        "iteration": next_iter,
        "mode": mode,
        "split": split,
        "total": total,
        "accuracy": round(accuracy, 4),
        "precision": round(precision, 4),
        "recall": round(recall, 4),
        "f1": round(f1, 4),
        "tp": tp, "tn": tn, "fp": fp, "fn": fn,
        "misclassified_count": len(misclassified),
        "added_to_vectorstore": added,
        "learned": learned,
        "note": improve_note,
        "eval_debug": {
            **eval_data,
            "mode": mode,
            "split": split,
        },
    }


@app.get("/rag/improve/history")
def rag_improve_history(user=Depends(require_role("admin"))):
    conn = get_db()
    cur = conn.cursor()
    try:
        cur.execute("SELECT * FROM improvement_history ORDER BY iteration ASC")
        rows = cur.fetchall()
    except Exception:
        rows = []
    conn.close()
    return [dict(r) for r in rows]


@app.delete("/rag/improve/history")
def rag_improve_history_clear(user=Depends(require_role("admin"))):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("DELETE FROM improvement_history")
    conn.commit()
    conn.close()
    return {"message": "?대젰 珥덇린???꾨즺"}


@app.post("/rag/reload")
def rag_reload(user=Depends(require_role("admin"))):
    """벡터스토어 디스크에서 강제 리로드"""
    import rag_engine, pickle, faiss, numpy as np
    vs_dir = rag_engine.VS_DIR
    index_path = os.path.join(vs_dir, "spam.index")
    meta_path = os.path.join(vs_dir, "metadata.pkl")
    try:
        with open(index_path, "rb") as f:
            new_index = faiss.deserialize_index(np.frombuffer(f.read(), dtype=np.uint8))
        with open(meta_path, "rb") as f:
            new_meta = pickle.load(f)
        rag_engine._index = new_index
        rag_engine._metadata = new_meta
        count = new_index.ntotal
        spam_cnt = new_meta["labels"].count("spam")
        ham_cnt = new_meta["labels"].count("ham")
        return {"ok": True, "vector_count": count, "spam_count": spam_cnt, "ham_count": ham_cnt, "path": vs_dir}
    except Exception as e:
        return {"ok": False, "error": str(e), "path": vs_dir}


@app.get("/rag/status")
def rag_status(user=Depends(require_role("admin"))):
    """RAG 踰≫꽣?ㅽ넗???꾪솴"""
    if not _rag_available:
        return {"available": False, "vector_count": 0}
    try:
        import pickle, faiss, numpy as np
        import rag_engine
        vs_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "vectorstore")
        idx_path = os.path.join(vs_dir, "spam.index")
        file_size = os.path.getsize(idx_path)
        with open(idx_path, "rb") as f:
            idx = faiss.deserialize_index(np.frombuffer(f.read(), dtype=np.uint8))
        with open(os.path.join(vs_dir, "metadata.pkl"), "rb") as f:
            meta = pickle.load(f)
        rag_engine._index = idx
        rag_engine._metadata = meta
        count = idx.ntotal
        spam_cnt = meta["labels"].count("spam")
        ham_cnt = meta["labels"].count("ham")
        return {"available": True, "vector_count": count, "spam_count": spam_cnt, "ham_count": ham_cnt}
    except Exception as e:
        return {"available": False, "error": str(e)}


# ?????????????????????????????
# GPT ?뚯씤?쒕떇
# ?????????????????????????????
@app.post("/rag/finetune/submit")
def finetune_submit(user=Depends(require_role("admin"))):
    """GPT ?뚯씤?쒕떇 ?묒뾽 ?쒖텧"""
    try:
        import finetune_gpt
        job_id = finetune_gpt.submit()
        return {"message": "?뚯씤?쒕떇 ?쒖텧 ?꾨즺", "job_id": job_id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/rag/finetune/status/{job_id}")
def finetune_status(job_id: str, user=Depends(require_role("admin"))):
    """?뚯씤?쒕떇 ?묒뾽 ?곹깭 ?뺤씤"""
    try:
        import finetune_gpt
        job = finetune_gpt.check_status(job_id)
        return {
            "job_id": job.id,
            "status": job.status,
            "model": job.model,
            "fine_tuned_model": job.fine_tuned_model,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/rag/finetune/jobs")
def finetune_jobs(user=Depends(require_role("admin"))):
    """?뚯씤?쒕떇 ?묒뾽 紐⑸줉"""
    try:
        import finetune_gpt
        jobs = finetune_gpt._load_jobs()
        return jobs
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/rag/finetune/apply/{job_id}")
def finetune_apply(job_id: str, user=Depends(require_role("admin"))):
    """?꾨즺???뚯씤?쒕떇 紐⑤뜽??GPT 遺꾨쪟???곸슜"""
    try:
        import finetune_gpt
        from rag_engine import classify_with_gpt
        job = finetune_gpt.check_status(job_id)
        if not job.fine_tuned_model:
            raise HTTPException(status_code=400, detail=f"?꾩쭅 ?꾨즺?섏? ?딆? ?묒뾽?낅땲?? ?꾩옱 ?곹깭: {job.status}")
        # classify_with_gpt??湲곕낯 紐⑤뜽???뚯씤?쒕떇 紐⑤뜽濡?援먯껜
        import rag_engine
        rag_engine._finetune_model = job.fine_tuned_model
        return {"message": f"?뚯씤?쒕떇 紐⑤뜽 ?곸슜 ?꾨즺: {job.fine_tuned_model}"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ?????????????????????????????
# ROOT
# ?????????????????????????????
@app.get("/langsmith/runs")
def langsmith_runs(limit: int = 50, user=Depends(require_role("developer"))):
    try:
        client = get_langsmith_client()
        if not client:
            raise HTTPException(status_code=503, detail="LangSmith 키가 설정되지 않았습니다.")
        project = os.getenv("LANGSMITH_PROJECT", "spam-detector")
        runs = list(client.list_runs(project_name=project, limit=limit))
        result = []
        for r in runs:
            result.append({
                "id": str(r.id),
                "name": r.name,
                "run_type": r.run_type,
                "status": r.status,
                "start_time": r.start_time.isoformat() if r.start_time else None,
                "end_time": r.end_time.isoformat() if r.end_time else None,
                "latency_ms": int((r.end_time - r.start_time).total_seconds() * 1000) if r.start_time and r.end_time else None,
                "total_tokens": getattr(r, "total_tokens", None),
                "prompt_tokens": getattr(r, "prompt_tokens", None),
                "completion_tokens": getattr(r, "completion_tokens", None),
                "inputs": r.inputs,
                "outputs": r.outputs,
                "error": r.error,
            })
        return result
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/langsmith/stats")
def langsmith_stats(user=Depends(require_role("developer"))):
    try:
        client = get_langsmith_client()
        if not client:
            raise HTTPException(status_code=503, detail="LangSmith 키가 설정되지 않았습니다.")
        project = os.getenv("LANGSMITH_PROJECT", "spam-detector")
        runs = list(client.list_runs(project_name=project, limit=100))
        total = len(runs)
        errors = sum(1 for r in runs if r.error)
        latencies = [(r.end_time - r.start_time).total_seconds() * 1000
                     for r in runs if r.start_time and r.end_time]
        avg_latency = int(sum(latencies) / len(latencies)) if latencies else 0
        total_tokens = sum(getattr(r, "total_tokens", 0) or 0 for r in runs)
        by_name: dict = {}
        for r in runs:
            by_name.setdefault(r.name, 0)
            by_name[r.name] += 1
        return {
            "total": total,
            "errors": errors,
            "avg_latency_ms": avg_latency,
            "total_tokens": total_tokens,
            "by_name": by_name,
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/")
def root():
    return {"status": "ok"}
